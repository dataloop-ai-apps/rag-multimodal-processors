"""
PPTX extraction logic.

Handles PowerPoint-specific extraction operations:
- Slide text extraction (titles, body, text boxes)
- Speaker notes extraction
- Image extraction with slide position metadata
- Table extraction
- Metadata collection
"""

import logging
import os
import tempfile
from typing import List, Tuple, Dict, Any

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from utils.extracted_data import ExtractedData
from utils.data_types import ImageContent

logger = logging.getLogger("rag-preprocessor")


class PPTXExtractor:
    """PPTX extraction operations."""

    @staticmethod
    def extract(data: ExtractedData) -> ExtractedData:
        """Extract content from PPTX item."""
        data.current_stage = "extraction"

        if not data.item:
            data.log_error("No item provided for extraction")
            return data

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                file_path = data.item.download(local_path=temp_dir)
                extract_images = data.config.extract_images
                extract_tables = data.config.extract_tables
                extract_notes = data.config.extract_notes

                content, images, metadata = PPTXExtractor._extract_pptx(
                    file_path, temp_dir, extract_images, extract_tables, extract_notes
                )

                data.content_text = content
                data.images = images
                metadata['source_file'] = data.item_name
                data.metadata = metadata

        except Exception:
            data.log_error("PPTX extraction failed. Check logs for details.")
            logger.exception("PPTX extraction error")

        return data

    @staticmethod
    def _extract_pptx(
        file_path: str,
        temp_dir: str,
        extract_images: bool,
        extract_tables: bool,
        extract_notes: bool = True,
    ) -> Tuple[str, List[ImageContent], Dict[str, Any]]:
        """Extract text, images and metadata from a PowerPoint file."""
        prs = Presentation(file_path)
        text_parts = []
        images = []
        table_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = PPTXExtractor._extract_slide_text(slide, slide_num)
            text_parts.append(slide_text)

            if extract_tables:
                table_text, n_tables = PPTXExtractor._extract_slide_tables(slide, slide_num)
                if table_text:
                    text_parts.append(table_text)
                table_count += n_tables

            if extract_notes:
                notes_text = PPTXExtractor._extract_slide_notes(slide, slide_num)
                if notes_text:
                    text_parts.append(notes_text)

            if extract_images:
                slide_images = PPTXExtractor._extract_slide_images(slide, slide_num, temp_dir)
                images.extend(slide_images)

        metadata = {
            'slide_count': len(prs.slides),
            'extraction_method': 'python-pptx',
            'image_count': len(images),
            'table_count': table_count,
            'processor': 'pptx',
        }

        return '\n'.join(text_parts), images, metadata

    @staticmethod
    def _extract_slide_text(slide: Any, slide_num: int) -> str:
        """Extract text from all text-bearing shapes in a slide."""
        parts = [f"\n\n--- Slide {slide_num} ---"]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                continue

            shape_text = []
            for para in shape.text_frame.paragraphs:
                para_text = ''.join(run.text for run in para.runs).strip()
                if para_text:
                    shape_text.append(para_text)

            if shape_text:
                parts.append('\n'.join(shape_text))

        return '\n'.join(parts)

    @staticmethod
    def _extract_slide_tables(slide: Any, slide_num: int) -> Tuple[str, int]:
        """Extract tables from a slide as formatted text."""
        table_texts = []
        table_count = 0

        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table_count += 1
            rows = []
            for row in shape.table.rows:
                row_text = ' | '.join(
                    cell.text.strip() for cell in row.cells
                )
                rows.append(row_text)
            if rows:
                table_texts.append(f"\n[Table on Slide {slide_num}]\n" + '\n'.join(rows))

        return '\n'.join(table_texts), table_count

    @staticmethod
    def _extract_slide_notes(slide: Any, slide_num: int) -> str:
        """Extract speaker notes from a slide."""
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip()
                if notes_text:
                    return f"\n[Speaker Notes - Slide {slide_num}]\n{notes_text}"
        except Exception:
            logger.debug("Could not extract notes from slide %d", slide_num)
        return ''

    @staticmethod
    def _extract_slide_images(slide: Any, slide_num: int, temp_dir: str) -> List[ImageContent]:
        """Extract images from a slide."""
        images = []

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                ext = image.ext
                image_path = os.path.join(temp_dir, f"slide{slide_num}_img{shape_idx}.{ext}")
                with open(image_path, 'wb') as f:
                    f.write(image.blob)

                left = shape.left / Inches(1) if shape.left is not None else None
                top = shape.top / Inches(1) if shape.top is not None else None
                width = shape.width / Inches(1) if shape.width is not None else None
                height = shape.height / Inches(1) if shape.height is not None else None
                bbox = (left, top, width, height) if all(v is not None for v in (left, top, width, height)) else None

                images.append(
                    ImageContent(
                        path=image_path,
                        page_number=slide_num,
                        format=ext,
                        size=(shape.width, shape.height),
                        bbox=bbox,
                    )
                )
            except (IOError, OSError, ValueError, AttributeError):
                logger.warning("Failed to extract image from slide %d shape %d", slide_num, shape_idx, exc_info=True)

        return images
